"""Conversions to/from PDF."""
import os
import shutil
import subprocess
import sys

import fitz

from .base import Job, JobError, open_doc, progress_cb, unique_path


# Get the results output directory from environment variable.
def _results_dir() -> str:
    return os.environ.get("PDF_STUDIO_RESULTS", ".")


# Extract all text from a PDF to a .txt file (plain or layout mode).
def pdf_to_text(inputs, options, job: Job):
    src = inputs[0]
    doc = open_doc(src)
    mode = options.get("mode", "plain")
    parts = []
    for i in range(doc.page_count):
        page = doc[i]
        parts.append(f"----- Page {i + 1} -----\n" + (page.get_text("text") or ""))
        progress_cb(job, i + 1, doc.page_count)
    doc.close()
    text = "\n\n".join(parts)
    name = os.path.splitext(os.path.basename(src))[0] + ".txt"
    path = unique_path(_results_dir(), name)
    with open(path, "w", encoding="utf-8") as f:
        f.write(text)
    return [(name, path)]


# Render each PDF page as an image (JPG/PNG) at specified DPI.
def _render_pages(inputs, options, job: Job, ext: str, default_dpi: int):
    src = inputs[0]
    doc = open_doc(src)
    try:
        dpi = int(options.get("dpi", default_dpi))
    except (TypeError, ValueError):
        dpi = default_dpi
    dpi = max(36, min(600, dpi))
    zoom = dpi / 72
    out = []
    base = os.path.splitext(os.path.basename(src))[0]
    for i in range(doc.page_count):
        pix = doc[i].get_pixmap(matrix=fitz.Matrix(zoom, zoom))
        name = f"{base}_{i + 1}.{ext}"
        path = unique_path(_results_dir(), name)
        pix.save(path)
        out.append((name, path))
        progress_cb(job, i + 1, doc.page_count)
    doc.close()
    return out


# Convert PDF pages to JPG images.
def pdf_to_jpg(inputs, options, job: Job):
    return _render_pages(inputs, options, job, "jpg", 150)


# Convert PDF pages to PNG images.
def pdf_to_png(inputs, options, job: Job):
    return _render_pages(inputs, options, job, "png", 150)


# Convert PDF to Word (.docx) using pdf2docx library.
def pdf_to_word(inputs, options, job: Job):
    src = inputs[0]
    try:
        from pdf2docx import Converter
    except ImportError:
        raise JobError("pdf2docx is not installed. Run: pip install pdf2docx")
    base = os.path.splitext(os.path.basename(src))[0]
    path = unique_path(_results_dir(), base + ".docx")
    cv = Converter(src)
    cv.convert(path, multi_processing=False)
    cv.close()
    return [(base + ".docx", path)]


# Extract tables from PDF to Excel (.xlsx) or CSV using pdfplumber.
def pdf_to_excel(inputs, options, job: Job):
    src = inputs[0]
    try:
        import pdfplumber
        from openpyxl import Workbook
    except ImportError:
        raise JobError("pdfplumber/openpyxl not installed.")
    base = os.path.splitext(os.path.basename(src))[0]
    fmt = options.get("format", "xlsx")
    ext = "csv" if fmt == "csv" else "xlsx"
    path = unique_path(_results_dir(), f"{base}.{ext}")
    wb = Workbook()
    ws = wb.active
    ws.title = "Tables"
    row0 = 1
    found = False
    with pdfplumber.open(src) as pdf:
        for pi, page in enumerate(pdf.pages):
            for table in page.extract_tables() or []:
                for r in table:
                    for c, cell in enumerate(r, start=1):
                        ws.cell(row=row0, column=c, value=(cell or "").strip() if cell else "")
                    row0 += 1
                row0 += 1
                found = True
            progress_cb(job, pi + 1, len(pdf.pages))
    if not found:
        raise JobError("No tables were detected in the document.")
    if fmt == "csv":
        import csv
        with open(path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            for row in ws.iter_rows(values_only=True):
                if any(v not in (None, "") for v in row):
                    writer.writerow(row)
    else:
        wb.save(path)
    return [(f"{base}.{ext}", path)]


# Convert PDF pages to a PowerPoint presentation (one slide per page).
def pdf_to_ppt(inputs, options, job: Job):
    src = inputs[0]
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise JobError("python-pptx not installed.")
    doc = open_doc(src)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    base = os.path.splitext(os.path.basename(src))[0]
    imgs = []
    try:
        for i in range(doc.page_count):
            pix = doc[i].get_pixmap(matrix=fitz.Matrix(2, 2))
            ip = unique_path(_results_dir(), f"_{base}_p{i + 1}.png")
            pix.save(ip)
            imgs.append(ip)
            progress_cb(job, i + 1, doc.page_count)
    finally:
        doc.close()
    for ip in imgs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(ip, 0, 0, width=prs.slide_width)
    name = base + ".pptx"
    path = unique_path(_results_dir(), name)
    prs.save(path)
    for ip in imgs:
        try:
            os.remove(ip)
        except OSError:
            pass
    return [(name, path)]


# Convert images (JPG/PNG/etc.) to a single PDF (one page per image).
def images_to_pdf(inputs, options, job: Job):
    doc = fitz.open()
    for i, path in enumerate(inputs):
        img = fitz.open(path)
        pix = fitz.Pixmap(img)
        page = doc.new_page(width=pix.width, height=pix.height)
        page.insert_image(page.rect, filename=path)
        img.close()
        progress_cb(job, i + 1, len(inputs))
    if doc.page_count == 0:
        raise JobError("No images provided.")
    name = options.get("name", "images") or "images"
    name = name if name.lower().endswith(".pdf") else name + ".pdf"
    path = unique_path(_results_dir(), name)
    doc.save(path, garbage=3, deflate=True)
    doc.close()
    return [(name, path)]


# Render HTML file to PDF using WeasyPrint (CSS-friendly).
def html_to_pdf(inputs, options, job: Job):
    src = inputs[0]
    try:
        from weasyprint import HTML
    except ImportError:
        raise JobError("weasyprint not installed.")
    with open(src, "r", encoding="utf-8", errors="replace") as f:
        html = f.read()
    name = os.path.splitext(os.path.basename(src))[0] + ".pdf"
    path = unique_path(_results_dir(), name)
    HTML(string=html, base_url=os.path.dirname(os.path.abspath(src))).write_pdf(path)
    return [(name, path)]


# LibreOffice headless conversion for docx/xlsx/pptx -> pdf.
def office_to_pdf(inputs, options, job: Job):
    src = inputs[0]
    soffice = options.get("soffice") or shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise JobError(
            "LibreOffice is not installed. Install it and add it to PATH, "
            "or enter the full path to soffice.exe."
        )
    outdir = unique_path(_results_dir(), "_lo_out")
    os.makedirs(outdir, exist_ok=True)
    cmd = [soffice, "--headless", "--convert-to", "pdf", "--outdir", outdir, src]
    try:
        proc = subprocess.run(cmd, capture_output=True, timeout=300)
        if proc.returncode != 0:
            raise JobError(f"LibreOffice failed: {proc.stderr.decode(errors='replace')[-500:]}")
    except subprocess.TimeoutExpired:
        raise JobError("LibreOffice conversion timed out.")
    outfile = os.path.join(outdir, os.path.splitext(os.path.basename(src))[0] + ".pdf")
    if not os.path.exists(outfile):
        found = [f for f in os.listdir(outdir) if f.lower().endswith(".pdf")]
        if not found:
            raise JobError("LibreOffice produced no output.")
        outfile = os.path.join(outdir, found[0])
    name = os.path.splitext(os.path.basename(src))[0] + ".pdf"
    final = unique_path(_results_dir(), name)
    shutil.move(outfile, final)
    shutil.rmtree(outdir, ignore_errors=True)
    return [(name, final)]


# Extract all embedded images from a PDF document.
def extract_images(inputs, options, job: Job):
    src = inputs[0]
    doc = open_doc(src)
    out = []
    img_no = 0
    total_xrefs = doc.xref_length()
    for xref in range(1, total_xrefs):
        if doc.xref_get_key(xref, "Subtype")[1] == "/Image":
            try:
                info = doc.extract_image(xref)
            except Exception:
                continue
            ext = info.get("ext", "png")
            if ext == "jpx":
                ext = "jp2"
            img_no += 1
            name = f"image_{img_no}.{ext}"
            path = unique_path(_results_dir(), name)
            with open(path, "wb") as f:
                f.write(info["image"])
            out.append((name, path))
            progress_cb(job, img_no, 50)  # loose progress
    doc.close()
    if not out:
        raise JobError("No images found in the document.")
    return out